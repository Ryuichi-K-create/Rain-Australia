# -*- coding: utf-8 -*-
"""
オーストラリア天気予測プレゼンテーション作成スクリプト
python-pptxを使用してPowerPointスライドを生成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# 定数
SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)

# カラーテーマ
THEME_BLUE = RGBColor(0, 112, 192)
THEME_DARK_BLUE = RGBColor(0, 51, 102)
THEME_LIGHT_BLUE = RGBColor(217, 226, 243)
THEME_GREEN = RGBColor(0, 128, 0)
THEME_RED = RGBColor(192, 0, 0)

# ベースパス
BASE_PATH = os.path.dirname(os.path.abspath(__file__))
FIGURES_PATH = os.path.join(BASE_PATH, "figures")


def add_title_slide(prs, title, subtitle):
    """タイトルスライドを追加"""
    slide_layout = prs.slide_layouts[6]  # 空白レイアウト
    slide = prs.slides.add_slide(slide_layout)

    # 背景色（グラデーション風に上部に色帯）
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME_BLUE
    shape.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = THEME_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # サブタイトル
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.3), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, section_title):
    """セクションタイトルスライドを追加"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SLIDE_WIDTH, Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME_BLUE
    shape.line.fill.background()

    # セクションタイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_func):
    """コンテンツスライドを追加"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # タイトルバー
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = THEME_BLUE
    title_bar.line.fill.background()

    # タイトルテキスト
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # コンテンツ追加（コールバック関数）
    content_func(slide)

    return slide


def add_text_with_bullets(slide, left, top, width, height, items, font_size=18, bold_first=False):
    """箇条書きテキストを追加"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = textbox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(6)
        if bold_first and i == 0:
            p.font.bold = True

    return textbox


def add_block(slide, left, top, width, title, content):
    """ブロック（色付き枠）を追加"""
    # ブロック背景
    block = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.2)
    )
    block.fill.solid()
    block.fill.fore_color.rgb = THEME_LIGHT_BLUE
    block.line.color.rgb = THEME_BLUE
    block.line.width = Pt(2)

    # ブロックテキスト
    textbox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = THEME_DARK_BLUE

    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(50, 50, 50)


def add_image_or_placeholder(slide, left, top, width, height, image_path, caption=""):
    """画像を追加（存在しない場合はプレースホルダー）"""
    full_path = os.path.join(BASE_PATH, image_path) if not os.path.isabs(image_path) else image_path

    if os.path.exists(full_path):
        slide.shapes.add_picture(full_path, Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        # プレースホルダー
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(180, 180, 180)
        placeholder.line.width = Pt(1)

        # プレースホルダーテキスト
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top + height/2 - 0.3), Inches(width), Inches(0.6))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[画像未生成]"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.alignment = PP_ALIGN.CENTER

    # キャプション
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(left), Inches(top + height + 0.05), Inches(width), Inches(0.4))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER


def add_table(slide, left, top, rows, cols, data, col_widths=None):
    """テーブルを追加"""
    table_width = sum(col_widths) if col_widths else Inches(cols * 1.5)
    row_height = Inches(0.4)

    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), table_width, row_height * rows)
    table = table_shape.table

    # 列幅設定
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # データ挿入
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)

            # ヘッダー行のスタイル
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME_BLUE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    paragraph.font.size = Pt(14)
                    paragraph.alignment = PP_ALIGN.CENTER
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.alignment = PP_ALIGN.CENTER

    return table_shape


def create_presentation():
    """メインのプレゼンテーション作成関数"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ===========================================
    # スライド1: タイトル
    # ===========================================
    add_title_slide(prs, "オーストラリア天気予測", "明日の雨予測モデルの構築と比較")

    # ===========================================
    # スライド2: 目次
    # ===========================================
    def content_toc(slide):
        sections = [
            "1. 課題概要",
            "2. データ分析（EDA）",
            "3. 前処理パイプライン",
            "4. モデル実装",
            "5. ハイパーパラメータチューニング",
            "6. 実験結果",
            "7. 考察",
            "8. まとめ"
        ]
        textbox = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(9), Inches(5.5))
        tf = textbox.text_frame
        for i, section in enumerate(sections):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = section
            p.font.size = Pt(26)
            p.font.color.rgb = THEME_DARK_BLUE
            p.space_after = Pt(16)

    add_content_slide(prs, "目次", content_toc)

    # ===========================================
    # スライド3: 課題の背景
    # ===========================================
    def content_background(slide):
        add_block(slide, 0.5, 1.4, 12.3, "目的",
                  "オーストラリアの気象データを用いて、明日が雨になるか否かを予測する2値分類モデルを構築する")

        items = [
            "データソース: Kaggle - Rain in Australia",
            "期間: 2007年11月〜2017年6月（約10年分）",
            "観測地点: オーストラリア全土49箇所",
            "応用: 農業計画、イベント運営、交通管理など"
        ]
        add_text_with_bullets(slide, 0.5, 2.9, 12, 4, items, font_size=22)

    add_content_slide(prs, "課題の背景", content_background)

    # ===========================================
    # スライド4: データセット概要
    # ===========================================
    def content_dataset(slide):
        # 左側テキスト
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.5), Inches(5))
        tf = textbox.text_frame

        p = tf.paragraphs[0]
        p.text = "基本情報"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        for item in ["レコード数: 145,460件", "特徴量数: 23カラム", "ターゲット: RainTomorrow (Yes/No)"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "クラス分布（不均衡）"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        for item in ["No（雨なし）: 78%", "Yes（雨あり）: 22%"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.space_after = Pt(4)

        # 右側画像
        add_image_or_placeholder(slide, 6.5, 1.5, 6, 4.5, "figures/target_distribution.png", "ターゲット変数の分布")

    add_content_slide(prs, "データセット概要", content_dataset)

    # ===========================================
    # スライド5: 特徴量の種類
    # ===========================================
    def content_features(slide):
        # 左列
        textbox1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(5.5))
        tf = textbox1.text_frame
        p = tf.paragraphs[0]
        p.text = "数値変数（16個）"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        num_features = [
            "気温: MinTemp, MaxTemp, Temp9am, Temp3pm",
            "降水: Rainfall, Evaporation",
            "日照: Sunshine",
            "風速: WindGustSpeed, WindSpeed9am/3pm",
            "湿度: Humidity9am/3pm",
            "気圧: Pressure9am/3pm",
            "雲量: Cloud9am/3pm"
        ]
        for item in num_features:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(4)

        # 右列
        textbox2 = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.5), Inches(5.5))
        tf = textbox2.text_frame
        p = tf.paragraphs[0]
        p.text = "カテゴリ変数（6個）"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        cat_features = [
            "Location: 49観測地点",
            "WindGustDir: 突風方向（16方位）",
            "WindDir9am: 9時の風向",
            "WindDir3pm: 15時の風向",
            "RainToday: 本日の雨（Yes/No）",
            "RainTomorrow: ターゲット"
        ]
        for item in cat_features:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(4)

    add_content_slide(prs, "特徴量の種類", content_features)

    # ===========================================
    # スライド6: 欠損値の分析
    # ===========================================
    def content_missing(slide):
        # 左側テキスト
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.5), Inches(5))
        tf = textbox.text_frame

        p = tf.paragraphs[0]
        p.text = "高欠損率の変数"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        for item in ["Sunshine: 48.0%", "Evaporation: 43.2%", "Cloud3pm: 40.8%", "Cloud9am: 38.4%"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "対処方針"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        for item in ["Location別の中央値で補完", "カテゴリ変数は最頻値で補完"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.space_after = Pt(4)

        # 右側画像
        add_image_or_placeholder(slide, 6.5, 1.5, 6, 4.5, "figures/missing_values.png", "欠損値ヒートマップ")

    add_content_slide(prs, "欠損値の分析", content_missing)

    # ===========================================
    # スライド7: データ分割
    # ===========================================
    def content_split(slide):
        add_block(slide, 0.5, 1.4, 12.3, "時系列分割の重要性",
                  "未来のデータで過去を予測することを防ぐため、時間順序を保持してデータを分割")

        # テーブル
        data = [
            ["データセット", "期間", "レコード数", "割合"],
            ["Train", "〜2015/06", "約113,000", "80%"],
            ["Validation", "2015/07〜2016/06", "約14,000", "10%"],
            ["Test", "2016/07〜", "約14,000", "10%"]
        ]
        col_widths = [Inches(2.5), Inches(3.5), Inches(2.5), Inches(1.5)]
        add_table(slide, 1.5, 3, 4, 4, data, col_widths)

    add_content_slide(prs, "データ分割（時系列）", content_split)

    # ===========================================
    # スライド8: 前処理の詳細
    # ===========================================
    def content_preprocess(slide):
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
        tf = textbox.text_frame

        sections = [
            ("1. 欠損値処理", ["数値変数: Location別の中央値で補完", "カテゴリ変数: Location別の最頻値で補完", "RainTomorrow欠損行は削除"]),
            ("2. 特徴量エンコーディング", ["風向（16方位）: サイクリカルエンコーディング（sinθ, cosθ）", "RainToday: バイナリ（Yes=1, No=0）", "Location: ラベルエンコーディング"]),
            ("3. 標準化", ["StandardScaler（平均0、標準偏差1）", "Trainデータで学習し、Val/Testに適用"])
        ]

        first = True
        for title, items in sections:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = title
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = THEME_DARK_BLUE
            p.space_before = Pt(10)

            for item in items:
                p = tf.add_paragraph()
                p.text = f"  • {item}"
                p.font.size = Pt(16)
                p.space_after = Pt(2)

    add_content_slide(prs, "前処理の詳細", content_preprocess)

    # ===========================================
    # スライド9: 使用モデル
    # ===========================================
    def content_models(slide):
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
        tf = textbox.text_frame

        models = [
            ("1. ロジスティック回帰（PyTorch実装）", ["線形モデル、解釈性が高い", "ベースラインとして使用"]),
            ("2. MLP（Multi-Layer Perceptron）（PyTorch実装）", ["非線形関係を学習可能", "BatchNorm + Dropoutで正則化"]),
            ("3. Random Forest（scikit-learn）", ["アンサンブル手法、過学習に強い", "特徴量重要度を算出可能"])
        ]

        first = True
        for title, items in models:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = title
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = THEME_DARK_BLUE
            p.space_before = Pt(16)

            for item in items:
                p = tf.add_paragraph()
                p.text = f"  • {item}"
                p.font.size = Pt(18)
                p.space_after = Pt(4)

    add_content_slide(prs, "使用モデル", content_models)

    # ===========================================
    # スライド10: クラス不均衡対策
    # ===========================================
    def content_imbalance(slide):
        add_block(slide, 0.5, 1.4, 12.3, "問題",
                  "Yes:No = 22:78 の不均衡データでは、「すべてNoと予測」でも78%の精度")

        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(12.3), Inches(4))
        tf = textbox.text_frame

        p = tf.paragraphs[0]
        p.text = "対策1: Weighted BCE Loss"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        p = tf.add_paragraph()
        p.text = "  L = -1/N Σ [w₁·y·log(ŷ) + w₀·(1-y)·log(1-ŷ)]"
        p.font.size = Pt(16)
        p.space_after = Pt(12)

        p = tf.add_paragraph()
        p.text = "対策2: Focal Loss"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        p = tf.add_paragraph()
        p.text = "  L_focal = -α(1-pₜ)^γ log(pₜ)"
        p.font.size = Pt(16)
        p.space_after = Pt(12)

        p = tf.add_paragraph()
        p.text = "対策3: class_weight='balanced'（Random Forest）"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

    add_content_slide(prs, "クラス不均衡対策", content_imbalance)

    # ===========================================
    # スライド11: Optunaによる最適化
    # ===========================================
    def content_optuna(slide):
        # 左列
        textbox1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(5.5))
        tf = textbox1.text_frame
        p = tf.paragraphs[0]
        p.text = "Optunaの特徴"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        for item in ["TPE（Tree-structured Parzen Estimator）", "自動枝刈り（Pruning）", "効率的な探索空間の探索"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "設定"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        for item in ["各モデル: 50トライアル", "最適化目標: Val AUC-ROC", "Pruner: MedianPruner"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(4)

        # 右列
        textbox2 = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.5), Inches(5.5))
        tf = textbox2.text_frame
        p = tf.paragraphs[0]
        p.text = "探索パラメータ"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        for item in ["学習率: [10⁻⁵, 10⁻¹]", "正則化: [10⁻⁶, 10⁻²]", "バッチサイズ: 64, 128, 256, 512", "損失関数: Weighted BCE, Focal", "オプティマイザ: SGD, Adam, AdamW"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(4)

    add_content_slide(prs, "Optunaによる最適化", content_optuna)

    # ===========================================
    # スライド12: ロジスティック回帰のチューニング
    # ===========================================
    def content_lr_tuning(slide):
        items = [
            "学習率: [10⁻⁴, 10⁻¹]（対数スケール）",
            "重み減衰（L2正則化）: [10⁻⁶, 10⁻²]",
            "オプティマイザ: SGD, Adam",
            "損失関数: Weighted BCE, Focal Loss"
        ]
        add_text_with_bullets(slide, 0.5, 1.4, 5.5, 2.5, items, font_size=18)
        add_image_or_placeholder(slide, 3, 3.8, 7, 3.2, "figures/learning_curves/logistic_regression.png", "ロジスティック回帰の学習曲線")

    add_content_slide(prs, "ロジスティック回帰のチューニング", content_lr_tuning)

    # ===========================================
    # スライド13: MLPのチューニング
    # ===========================================
    def content_mlp_tuning(slide):
        items = [
            "隠れ層数: 1〜4層",
            "各層のユニット数: 32〜512",
            "Dropout率: 0.0〜0.5",
            "オプティマイザ: Adam, AdamW"
        ]
        add_text_with_bullets(slide, 0.5, 1.4, 5.5, 2.5, items, font_size=18)
        add_image_or_placeholder(slide, 3, 3.8, 7, 3.2, "figures/learning_curves/mlp.png", "MLPの学習曲線")

    add_content_slide(prs, "MLPのチューニング", content_mlp_tuning)

    # ===========================================
    # スライド14: Random Forestのチューニング
    # ===========================================
    def content_rf_tuning(slide):
        items = [
            "n_estimators: 50〜500",
            "max_depth: 5〜50",
            "min_samples_split: 2〜20",
            "min_samples_leaf: 1〜10",
            "max_features: 'sqrt', 'log2', None"
        ]
        add_text_with_bullets(slide, 0.5, 1.4, 5.5, 2.5, items, font_size=18)
        add_image_or_placeholder(slide, 3, 3.8, 7, 3.2, "figures/feature_importance/random_forest.png", "Random Forestの特徴量重要度（上位10）")

    add_content_slide(prs, "Random Forestのチューニング", content_rf_tuning)

    # ===========================================
    # スライド15: 評価指標
    # ===========================================
    def content_metrics(slide):
        data = [
            ["指標", "説明"],
            ["Accuracy", "全体の正解率"],
            ["Precision", "雨と予測した中で実際に雨だった割合"],
            ["Recall", "実際の雨をどれだけ捕捉できたか"],
            ["F1-Score", "PrecisionとRecallの調和平均"],
            ["AUC-ROC", "ROC曲線下面積（主指標）"],
            ["AUC-PR", "Precision-Recall曲線下面積"]
        ]
        col_widths = [Inches(2.5), Inches(8)]
        add_table(slide, 1.2, 1.5, 7, 2, data, col_widths)

        add_block(slide, 1.2, 5.0, 10.5, "AUC-ROCを主指標とする理由",
                  "閾値に依存しない評価が可能であり、クラス不均衡データに対して適切")

    add_content_slide(prs, "評価指標", content_metrics)

    # ===========================================
    # スライド16: テストデータでの評価結果
    # ===========================================
    def content_results(slide):
        data = [
            ["モデル", "Accuracy", "Precision", "Recall", "F1", "AUC-ROC", "AUC-PR"],
            ["Logistic Reg.", "0.785", "0.526", "0.742", "0.616", "0.852", "0.674"],
            ["MLP", "0.797", "0.542", "0.785", "0.641", "0.875", "0.720"],
            ["Random Forest", "0.843", "0.764", "0.463", "0.577", "0.873", "0.715"]
        ]
        col_widths = [Inches(2.2), Inches(1.5), Inches(1.5), Inches(1.2), Inches(1), Inches(1.5), Inches(1.3)]
        add_table(slide, 0.8, 1.5, 4, 7, data, col_widths)

        add_image_or_placeholder(slide, 3.5, 3.8, 6, 3, "figures/model_comparison.png", "評価指標の比較")

    add_content_slide(prs, "テストデータでの評価結果", content_results)

    # ===========================================
    # スライド17: ROCカーブ比較
    # ===========================================
    def content_roc(slide):
        add_image_or_placeholder(slide, 2.5, 1.4, 8.5, 4.5, "figures/roc_curves/comparison.png", "ROCカーブの比較")

        items = ["曲線が左上に近いほど性能が良い", "対角線はランダム予測（AUC=0.5）"]
        add_text_with_bullets(slide, 0.5, 6.2, 12, 1, items, font_size=16)

    add_content_slide(prs, "ROCカーブ比較", content_roc)

    # ===========================================
    # スライド18: 混同行列
    # ===========================================
    def content_confusion(slide):
        add_image_or_placeholder(slide, 1.5, 1.4, 10.5, 5.5, "figures/confusion_matrices/comparison.png", "各モデルの混同行列")

    add_content_slide(prs, "混同行列", content_confusion)

    # ===========================================
    # スライド19: モデル比較の考察
    # ===========================================
    def content_discussion(slide):
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
        tf = textbox.text_frame

        sections = [
            ("Random Forestが最高性能を達成", ["非線形関係を効果的に捉える", "特徴量間の相互作用を自動学習", "アンサンブル効果で過学習を抑制"]),
            ("MLPの特徴", ["良好な性能だがチューニングに時間がかかる", "適切なアーキテクチャ選択が重要", "GPU活用で学習を高速化"]),
            ("ロジスティック回帰の特徴", ["シンプルで解釈性が高い", "非線形関係の捕捉に限界", "ベースラインとして有用"])
        ]

        first = True
        for title, items in sections:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = THEME_DARK_BLUE
            p.space_before = Pt(10)

            for item in items:
                p = tf.add_paragraph()
                p.text = f"  • {item}"
                p.font.size = Pt(16)
                p.space_after = Pt(2)

    add_content_slide(prs, "モデル比較の考察", content_discussion)

    # ===========================================
    # スライド20: 重要な特徴量
    # ===========================================
    def content_importance(slide):
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(2))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "Random Forestの特徴量重要度分析から:"
        p.font.size = Pt(18)

        textbox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.3), Inches(3))
        tf = textbox2.text_frame
        p = tf.paragraphs[0]
        p.text = "上位の特徴量"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        features = [
            "1. Humidity3pm: 午後の湿度が最重要",
            "2. Pressure3pm: 気圧も重要な予測因子",
            "3. Sunshine: 日照時間",
            "4. Cloud3pm: 午後の雲量",
            "5. RainToday: 本日の雨の有無"
        ]
        for item in features:
            p = tf.add_paragraph()
            p.text = f"  {item}"
            p.font.size = Pt(18)
            p.space_after = Pt(4)

        add_block(slide, 0.5, 5.5, 12.3, "気象学的解釈",
                  "午後の湿度・気圧・雲量は、翌日の降雨を予測する上で物理的に妥当な指標")

    add_content_slide(prs, "重要な特徴量", content_importance)

    # ===========================================
    # スライド21: まとめ
    # ===========================================
    def content_summary(slide):
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
        tf = textbox.text_frame

        p = tf.paragraphs[0]
        p.text = "本研究の成果"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        achievements = [
            "3種類のモデル（ロジスティック回帰、MLP、Random Forest）を実装・比較",
            "Optunaによる体系的なハイパーパラメータチューニングを実施",
            "クラス不均衡に対する適切な対策（Weighted Loss, class_weight）を適用",
            "MLPが最高のAUC-ROC(0.875)を達成"
        ]
        for i, item in enumerate(achievements):
            p = tf.add_paragraph()
            p.text = f"  {i+1}. {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "今後の課題"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_BLUE

        future = [
            "勾配ブースティング（LightGBM、XGBoost）の追加",
            "時系列特徴量（ラグ特徴量、移動平均）の導入",
            "モデルアンサンブル（スタッキング）の検討",
            "より多くのトライアルでのハイパーパラメータ探索"
        ]
        for item in future:
            p = tf.add_paragraph()
            p.text = f"  • {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(4)

    add_content_slide(prs, "まとめ", content_summary)

    # ===========================================
    # スライド22: 参考文献
    # ===========================================
    def content_references(slide):
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
        tf = textbox.text_frame

        references = [
            "[1] Young, J. (2017). Rain in Australia. Kaggle Dataset.\n     https://www.kaggle.com/datasets/jsphyg/weather-dataset-rattle-package",
            "[2] Paszke, A., et al. (2019). PyTorch: An Imperative Style,\n     High-Performance Deep Learning Library. NeurIPS.",
            "[3] Akiba, T., et al. (2019). Optuna: A Next-generation\n     Hyperparameter Optimization Framework. KDD.",
            "[4] Pedregosa, F., et al. (2011). Scikit-learn: Machine Learning\n     in Python. JMLR."
        ]

        for i, ref in enumerate(references):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = ref
            p.font.size = Pt(16)
            p.space_after = Pt(16)

    add_content_slide(prs, "参考文献", content_references)

    # 保存
    output_path = os.path.join(BASE_PATH, "presentation.pptx")
    prs.save(output_path)
    print(f"プレゼンテーションを保存しました: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
