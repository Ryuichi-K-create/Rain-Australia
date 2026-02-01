# -*- coding: utf-8 -*-
"""
オーストラリア天気予測プレゼンテーション作成スクリプト
python-pptxを使用してPowerPointスライドを生成
presentation.texの内容に基づいて更新
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 定数
SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)

# カラーテーマ（黒基調）
THEME_PRIMARY = RGBColor(40, 40, 40)       # メイン（濃いグレー/黒）
THEME_DARK = RGBColor(20, 20, 20)          # 濃い黒
THEME_LIGHT = RGBColor(230, 230, 230)      # 薄いグレー（ブロック背景）
THEME_ACCENT = RGBColor(80, 80, 80)        # アクセント（中間グレー）

# フォント設定
FONT_JAPANESE = "游ゴシック"  # 日本語用ゴシック体
FONT_ENGLISH = "Times New Roman"  # 英語用

# ベースパス
BASE_PATH = os.path.dirname(os.path.abspath(__file__))
FIGURES_PATH = os.path.join(BASE_PATH, "figures")


def set_font(run, text, size, bold=False, color=None):
    """フォントを設定（日本語はゴシック、英語はTimes New Roman）"""
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    # 英語フォント設定
    run.font.name = FONT_ENGLISH
    # 日本語フォント設定（East Asian font）
    from pptx.oxml.ns import qn
    from lxml import etree
    rPr = run.font._element
    # 東アジアフォントを設定
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', FONT_JAPANESE)


def add_title_slide(prs, title, subtitle):
    """タイトルスライドを追加"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景色帯
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME_PRIMARY
    shape.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, title, 52, bold=True, color=THEME_DARK)
    p.alignment = PP_ALIGN.CENTER

    # サブタイトル
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.3), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, subtitle, 32, color=RGBColor(80, 80, 80))
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title):
    """コンテンツスライドを追加（タイトルバー付き）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # タイトルバー
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = THEME_PRIMARY
    title_bar.line.fill.background()

    # タイトルテキスト
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, title, 36, bold=True, color=RGBColor(255, 255, 255))

    return slide


def add_block(slide, left, top, width, height, title, content):
    """ブロック（色付き枠）を追加"""
    block = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    block.fill.solid()
    block.fill.fore_color.rgb = THEME_LIGHT
    block.line.color.rgb = THEME_PRIMARY
    block.line.width = Pt(2)

    textbox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(height - 0.3))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, title, 20, bold=True, color=THEME_DARK)

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    set_font(run2, content, 18, color=RGBColor(50, 50, 50))


def add_text(slide, left, top, width, height, text, size=20, bold=False, color=None):
    """テキストボックスを追加"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, text, size, bold, color or RGBColor(50, 50, 50))
    return tf


def add_bullets(slide, left, top, width, height, title, items, title_size=24, item_size=20):
    """タイトル付き箇条書きを追加"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = textbox.text_frame
    tf.word_wrap = True

    # タイトル
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, title, title_size, bold=True, color=THEME_DARK)
    p.space_after = Pt(8)

    # 箇条書き
    for item in items:
        p = tf.add_paragraph()
        run = p.add_run()
        set_font(run, f"・{item}", item_size, color=RGBColor(50, 50, 50))
        p.space_after = Pt(6)

    return tf


def add_image(slide, left, top, max_width, max_height, image_path, caption=""):
    """画像を追加（アスペクト比を維持、存在しない場合はプレースホルダー）"""
    full_path = os.path.join(FIGURES_PATH, image_path)

    if os.path.exists(full_path):
        from PIL import Image
        # 画像の元サイズを取得
        with Image.open(full_path) as img:
            img_width, img_height = img.size

        # アスペクト比を計算
        aspect_ratio = img_width / img_height
        max_aspect = max_width / max_height

        # 指定領域内でアスペクト比を維持した最大サイズを計算
        if aspect_ratio > max_aspect:
            # 横長の画像：幅に合わせる
            width = max_width
            height = max_width / aspect_ratio
        else:
            # 縦長の画像：高さに合わせる
            height = max_height
            width = max_height * aspect_ratio

        # 中央揃えのためのオフセット計算
        left_offset = (max_width - width) / 2
        top_offset = (max_height - height) / 2

        slide.shapes.add_picture(full_path, Inches(left + left_offset), Inches(top + top_offset),
                                  width=Inches(width), height=Inches(height))
    else:
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(180, 180, 180)

        textbox = slide.shapes.add_textbox(Inches(left), Inches(top + height/2 - 0.3), Inches(width), Inches(0.6))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        set_font(run, "[画像未生成]", 16, color=RGBColor(150, 150, 150))
        p.alignment = PP_ALIGN.CENTER

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(left), Inches(top + height + 0.1), Inches(width), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        set_font(run, caption, 14, color=RGBColor(100, 100, 100))
        p.alignment = PP_ALIGN.CENTER


def add_table(slide, left, top, data, col_widths):
    """テーブルを追加"""
    rows = len(data)
    cols = len(data[0])
    table_width = sum(col_widths)
    row_height = Inches(0.45)

    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), table_width, row_height * rows)
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME_PRIMARY
                set_font(run, str(cell_text), 16, bold=True, color=RGBColor(255, 255, 255))
            else:
                set_font(run, str(cell_text), 14, color=RGBColor(50, 50, 50))
            p.alignment = PP_ALIGN.CENTER

    return table_shape


def create_presentation():
    """メインのプレゼンテーション作成関数"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ===========================================
    # スライド1: タイトル
    # ===========================================
    add_title_slide(prs, "天気予測", "機械科学専攻 1年\n小西隆一")

    # ===========================================
    # スライド2: 課題の背景
    # ===========================================
    slide = add_content_slide(prs, "課題の背景")
    add_block(slide, 0.5, 1.5, 12.3, 1.2, "目的",
              "オーストラリアの気象データを用いて、明日が雨になるか否かを予測する2値分類モデルを構築する")

    add_bullets(slide, 0.5, 3.0, 6, 2, "なぜこのテーマを選んだか",
                ["金沢市は年間降水量が多く、雨の日を事前に予測できると日常生活で役立つ",
                 "気象データを用いた機械学習予測に興味があった",
                 "2値分類の実践的な題材として適切だった"], 22, 18)

    add_bullets(slide, 6.8, 3.0, 6, 2, "日常生活・研究への応用",
                ["外出・イベントの計画立案",
                 "農業における作業スケジュール管理",
                 "交通機関の運行計画への活用"], 22, 18)

    # ===========================================
    # スライド3: データセット概要
    # ===========================================
    slide = add_content_slide(prs, "データセット概要")

    add_bullets(slide, 0.5, 1.5, 5.5, 2.5, "基本情報",
                ["レコード数: 145,460件",
                 "特徴量数: 23カラム",
                 "ターゲット: RainTomorrow (Yes/No)",
                 "期間: 2007年11月〜2017年6月",
                 "観測地点: オーストラリア全土49箇所"], 22, 18)

    add_bullets(slide, 0.5, 4.5, 5.5, 1.5, "クラス分布（不均衡）",
                ["No（雨なし）: 78%",
                 "Yes（雨あり）: 22%"], 22, 18)

    add_image(slide, 6.5, 1.5, 6, 4.5, "target_distribution.png", "ターゲット変数の分布")

    # ===========================================
    # スライド4: 特徴量の種類
    # ===========================================
    slide = add_content_slide(prs, "特徴量の種類")

    add_bullets(slide, 0.5, 1.5, 6, 5.5, "数値変数（16個）",
                ["気温: MinTemp, MaxTemp, Temp9am, Temp3pm",
                 "降水: Rainfall, Evaporation",
                 "日照: Sunshine",
                 "風速: WindGustSpeed, WindSpeed9am/3pm",
                 "湿度: Humidity9am/3pm",
                 "気圧: Pressure9am/3pm",
                 "雲量: Cloud9am/3pm"], 24, 18)

    add_bullets(slide, 6.8, 1.5, 6, 5.5, "カテゴリ変数（6個）",
                ["Location: 49観測地点",
                 "WindGustDir: 突風方向（16方位）",
                 "WindDir9am: 9時の風向",
                 "WindDir3pm: 15時の風向",
                 "RainToday: 本日の雨（Yes/No）",
                 "RainTomorrow: ターゲット"], 24, 18)

    # ===========================================
    # スライド5: データ分割（時系列）
    # ===========================================
    slide = add_content_slide(prs, "データ分割（時系列）")
    add_block(slide, 0.5, 1.5, 12.3, 1.2, "時系列分割の重要性",
              "未来のデータで過去を予測することを防ぐため、時間順序を保持してデータを分割")

    data = [
        ["データセット", "期間", "レコード数", "割合"],
        ["Train", "〜2015/06", "約113,000", "80%"],
        ["Validation", "2015/07〜2016/06", "約14,000", "10%"],
        ["Test", "2016/07〜", "約14,000", "10%"]
    ]
    col_widths = [Inches(3), Inches(3.5), Inches(2.5), Inches(1.5)]
    add_table(slide, 1.5, 3.2, data, col_widths)

    # ===========================================
    # スライド6: 前処理の詳細
    # ===========================================
    slide = add_content_slide(prs, "前処理の詳細")

    add_bullets(slide, 0.5, 1.5, 12, 1.8, "1. 欠損値処理",
                ["数値変数: Location別の中央値で補完",
                 "カテゴリ変数: Location別の最頻値で補完",
                 "RainTomorrow欠損行は削除"], 22, 18)

    add_bullets(slide, 0.5, 3.5, 12, 1.8, "2. 特徴量エンコーディング",
                ["風向（16方位）: サイクリカルエンコーディング（sin θ, cos θ）",
                 "RainToday: バイナリ（Yes=1, No=0）",
                 "Location: ラベルエンコーディング"], 22, 18)

    add_bullets(slide, 0.5, 5.5, 12, 1.5, "3. 標準化",
                ["StandardScaler（平均0、標準偏差1）",
                 "Trainデータで学習し、Val/Testに適用"], 22, 18)

    # ===========================================
    # スライド7: 使用モデル
    # ===========================================
    slide = add_content_slide(prs, "使用モデル")

    add_bullets(slide, 0.5, 1.5, 12, 1.8, "1. ロジスティック回帰",
                ["構造: 入力層 → 出力層（1ユニット）のシンプルな線形モデル",
                 "特長: 解釈性が高く、各特徴量の重みから予測への寄与度がわかる"], 22, 18)

    add_bullets(slide, 0.5, 3.5, 12, 1.8, "2. MLP（Multi-Layer Perceptron）",
                ["構造: 入力層 → 隠れ層（複数）→ 出力層の多層ニューラルネット",
                 "特長: 非線形な関係を学習可能、BatchNormとDropoutで過学習を抑制"], 22, 18)

    add_bullets(slide, 0.5, 5.5, 12, 1.5, "3. Random Forest",
                ["構造: 複数の決定木を並列に学習し、多数決で予測するアンサンブル手法",
                 "特長: 過学習に強く、特徴量重要度を算出可能"], 22, 18)

    # ===========================================
    # スライド8: クラス不均衡対策
    # ===========================================
    slide = add_content_slide(prs, "クラス不均衡対策")
    add_block(slide, 0.5, 1.5, 12.3, 1.0, "問題点",
              "Yes:No = 22:78 の不均衡データでは、「すべてNoと予測」しても78%の精度になってしまう")

    add_text(slide, 0.5, 2.8, 12, 0.5, "実施した対策", 24, bold=True, color=THEME_DARK)

    add_bullets(slide, 0.5, 3.4, 12, 1.8, "1. ロジスティック回帰・MLP",
                ["損失関数に重み付けを導入（Weighted BCE Loss）",
                 "少数クラス（雨あり）の誤分類に対するペナルティを大きく設定"], 20, 18)

    add_bullets(slide, 0.5, 5.3, 12, 1.5, "2. Random Forest",
                ["class_weight='balanced' パラメータを使用",
                 "各クラスのサンプル数に応じて自動的に重みを調整"], 20, 18)

    # ===========================================
    # スライド9: ロジスティック回帰のチューニング
    # ===========================================
    slide = add_content_slide(prs, "ロジスティック回帰のチューニング")

    add_bullets(slide, 0.5, 1.5, 5.5, 2.2, "探索したハイパーパラメータ",
                ["学習率: モデルの重み更新の大きさを制御",
                 "重み減衰（L2正則化）: 過学習を防ぐ",
                 "バッチサイズ: 128, 256, 512"], 20, 16)

    add_block(slide, 0.5, 4.0, 5.5, 2.5, "選定されたパラメータ",
              "・学習率: 0.0135\n・重み減衰: 4.15×10⁻⁶\n・バッチサイズ: 512\n")

    add_image(slide, 6.3, 1.5, 6.5, 5, "learning_curves/logistic_regression.png", "学習曲線")

    # ===========================================
    # スライド10: MLPのチューニング
    # ===========================================
    slide = add_content_slide(prs, "MLPのチューニング")

    add_bullets(slide, 0.5, 1.5, 5.5, 2.2, "探索したハイパーパラメータ",
                ["隠れ層数: 1〜3層",
                 "各層のユニット数: 64〜256",
                 "Dropout率: 0.1〜0.4"], 20, 16)

    add_block(slide, 0.5, 4.0, 5.5, 2.5, "選定されたパラメータ",
              "・隠れ層: 2層 [64, 192]\n・Dropout率: 0.163\n・学習率: 0.00168\n・バッチサイズ: 256\n")

    add_image(slide, 6.3, 1.5, 6.5, 5, "learning_curves/mlp.png", "学習曲線")

    # ===========================================
    # スライド11: Random Forestのチューニング
    # ===========================================
    slide = add_content_slide(prs, "Random Forestのチューニング")

    add_bullets(slide, 0.5, 1.5, 6, 3.5, "探索したハイパーパラメータ",
                ["n_estimators: 決定木の数（100〜300）",
                 "max_depth: 各決定木の最大深さ（10〜30）",
                 "class_weight: 'balanced'（固定）"], 20, 16)

    add_block(slide, 0.5, 5.2, 6, 1.5, "選定されたパラメータ",
              "・n_estimators: 300\n・max_depth: 26\n")

    add_image(slide, 6.8, 1.5, 6, 5, "feature_importance/random_forest.png", "特徴量重要度（上位15）")

    # ===========================================
    # スライド12: 評価指標
    # ===========================================
    slide = add_content_slide(prs, "評価指標")

    data = [
        ["指標", "説明"],
        ["Accuracy", "全体の正解率。全予測のうち正しく分類できた割合"],
        ["Precision", "「雨」と予測したもののうち、実際に雨だった割合"],
        ["Recall", "実際の雨の日のうち、正しく「雨」と予測できた割合"],
        ["F1-Score", "PrecisionとRecallの調和平均"]
    ]
    col_widths = [Inches(2.5), Inches(9)]
    add_table(slide, 0.8, 1.5, data, col_widths)

    add_block(slide, 0.8, 5.0, 11.5, 1.3, "不均衡データでの注意点",
              "Accuracyだけでは正しく評価できない。Precision, Recall, F1-Scoreを総合的に見ることが重要")

    # ===========================================
    # スライド13: テストデータでの評価結果
    # ===========================================
    slide = add_content_slide(prs, "テストデータでの評価結果")

    data = [
        ["モデル", "Accuracy", "Precision", "Recall", "F1-Score"],
        ["ロジスティック回帰", "0.785", "0.526", "0.742", "0.616"],
        ["MLP", "0.797", "0.542", "0.785", "0.641"],
        ["Random Forest", "0.843", "0.764", "0.463", "0.577"]
    ]
    col_widths = [Inches(3), Inches(2), Inches(2), Inches(2), Inches(2)]
    add_table(slide, 1, 1.5, data, col_widths)

    add_bullets(slide, 0.5, 4.2, 12, 2.5, "結果の概要",
                ["Random ForestがAccuracyとPrecisionで最高",
                 "MLPがRecallとF1-Scoreで最高",
                 "ロジスティック回帰はベースラインとして妥当な性能"], 22, 20)

    # ===========================================
    # スライド14: 混同行列
    # ===========================================
    slide = add_content_slide(prs, "混同行列")
    add_image(slide, 1.5, 1.4, 10.5, 5.5, "confusion_matrices/comparison.png", "各モデルの混同行列")

    # ===========================================
    # スライド15: モデル比較の考察
    # ===========================================
    slide = add_content_slide(prs, "モデル比較の考察")

    add_bullets(slide, 0.5, 1.5, 12, 1.6, "Random Forestの高いAccuracy・Precisionの理由",
                ["決定木のアンサンブルにより、特徴量間の複雑な相互作用を捉えられた",
                 "ただしRecallが低く、雨の日の見逃しが多い"], 20, 18)

    add_bullets(slide, 0.5, 3.3, 12, 1.6, "MLPの高いRecall・F1-Scoreの理由",
                ["重み付き損失関数により、少数クラス（雨あり）を積極的に検出",
                 "Precisionは低めで、誤警報が多い傾向"], 20, 18)

    add_bullets(slide, 0.5, 5.1, 12, 1.6, "ロジスティック回帰の限界",
                ["線形モデルのため、非線形な関係の捕捉に限界がある",
                 "しかしベースラインとして、他モデルの改善度を測る基準になった"], 20, 18)

    # ===========================================
    # スライド16: 重要な特徴量
    # ===========================================
    slide = add_content_slide(prs, "重要な特徴量")
    add_text(slide, 0.5, 1.5, 6, 0.5, "Random Forestの特徴量重要度分析から:", 20)

    add_bullets(slide, 0.5, 2.1, 5.5, 3, "上位の特徴量",
                ["Humidity3pm: 午後の湿度",
                 "Pressure3pm: 午後の気圧",
                 "Sunshine: 日照時間",
                 "Cloud3pm: 午後の雲量",
                 "RainToday: 本日の雨の有無"], 22, 18)

    add_block(slide, 0.5, 5.5, 5.5, 1.2, "気象学的解釈",
              "午後の湿度・気圧・雲量は、翌日の降雨を予測する上で物理的に妥当な指標")

    add_image(slide, 6.5, 1.5, 6.3, 5, "feature_importance/random_forest.png", "特徴量重要度（上位10）")

    # ===========================================
    # スライド17: まとめ
    # ===========================================
    slide = add_content_slide(prs, "まとめ")

    add_bullets(slide, 0.5, 1.5, 12, 3.2, "本研究の成果",
                ["3種類のモデル（ロジスティック回帰、MLP、Random Forest）を実装・比較",
                 "体系的なハイパーパラメータチューニングを実施",
                 "クラス不均衡に対する適切な対策（重み付き損失関数）を適用",
                 "Random ForestがAccuracy 0.843、Precision 0.764で最高性能を達成"], 24, 20)

    add_bullets(slide, 0.5, 4.9, 12, 2, "今後の課題",
                ["勾配ブースティング（LightGBM、XGBoost）の追加検討",
                 "時系列特徴量（ラグ特徴量、移動平均）の導入",
                 "Recallを重視したモデル調整（雨の見逃しを減らす）"], 24, 20)

    # 保存
    output_path = os.path.join(BASE_PATH, "presentation.pptx")
    prs.save(output_path)
    print(f"プレゼンテーションを保存しました: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
